import uuid
from pathlib import Path
import pypdf
import csv


def _chunk_text(text: str, source: str, file_path: str) -> list[dict]:
    chunks = []
    chunk_size, overlap = 1000, 200
    for i in range(0, len(text), chunk_size - overlap):
        chunk = text[i : i + chunk_size].strip()
        if len(chunk) < 50:
            continue
        chunks.append({
            "id": str(uuid.uuid4()),
            "text": chunk,
            "metadata": {"source": source, "file_path": file_path, "chunk_index": len(chunks)},
        })
    return chunks


def load_pdf(file_path: str) -> list[dict]:
    reader = pypdf.PdfReader(file_path)
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    return _chunk_text(text, Path(file_path).stem, file_path)


def load_txt(file_path: str) -> list[dict]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return _chunk_text(text, Path(file_path).stem, file_path)


def load_csv(file_path: str) -> list[dict]:
    rows = []
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        reader = csv.DictReader(f)
        for row in reader:
            rows.append(", ".join(f"{k}: {v}" for k, v in row.items()))
    text = "\n".join(rows)
    return _chunk_text(text, Path(file_path).stem, file_path)


def load_docx(file_path: str) -> list[dict]:
    from docx import Document
    doc = Document(file_path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return _chunk_text(text, Path(file_path).stem, file_path)


def load_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    text = "\n".join(lines)
    return _chunk_text(text, Path(file_path).stem, file_path)


def load_xlsx(file_path: str) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = ", ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                lines.append(row_text)
    text = "\n".join(lines)
    return _chunk_text(text, Path(file_path).stem, file_path)


LOADERS = {
    ".pdf": load_pdf,
    ".txt": load_txt,
    ".csv": load_csv,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


def load_documents_from_dir(directory: str) -> list[dict]:
    all_chunks = []
    for file in Path(directory).rglob("*"):
        loader = LOADERS.get(file.suffix.lower())
        if not loader:
            continue
        print(f"Loading: {file.name}")
        try:
            chunks = loader(str(file))
            all_chunks.extend(chunks)
            print(f"  → {len(chunks)} chunks")
        except Exception as e:
            print(f"  ⚠ 오류: {e}")
    print(f"\n총 {len(all_chunks)}개 청크 로드 완료")
    return all_chunks
